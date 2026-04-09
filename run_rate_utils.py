import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from io import BytesIO
import warnings
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime, timedelta, date

# ==============================================================================
# --- 1. CONSTANTS & UTILITY FUNCTIONS ---
# ==============================================================================

PASTEL_COLORS = {
    'red': '#ff6961',
    'orange': '#ffb347',
    'green': '#77dd77',
    'blue': '#3498DB'
}


def format_minutes_to_dhm(total_minutes):
    """Converts total minutes into a human-readable duration string."""
    if pd.isna(total_minutes) or total_minutes < 0:
        return "N/A"
    if total_minutes < 1.0:
        return f"{total_minutes * 60:.0f}s"
    total_seconds = total_minutes * 60
    total_mins_int = int(total_seconds // 60)
    remaining_secs = int(total_seconds % 60)
    if total_mins_int < 10:
        return f"{total_mins_int}m {remaining_secs:02d}s" if remaining_secs else f"{total_mins_int}m"
    days = total_mins_int // (60 * 24)
    remaining_minutes = total_mins_int % (60 * 24)
    hours = remaining_minutes // 60
    minutes = remaining_minutes % 60
    parts = []
    if days > 0: parts.append(f"{days}d")
    if hours > 0: parts.append(f"{hours}h")
    if minutes > 0 or not parts: parts.append(f"{minutes}m")
    return " ".join(parts) if parts else "0m"


def format_duration(seconds):
    """Converts total seconds into a 'Xd Yh Zm' string."""
    if pd.isna(seconds) or seconds < 0:
        return "N/A"
    return format_minutes_to_dhm(seconds / 60)


def get_renamed_summary_df(df_in):
    """Helper function to rename summary tables consistently."""
    if df_in is None or df_in.empty:
        return pd.DataFrame()
    df = df_in.copy()

    rename_map = {
        'hour': 'Hour', 'date': 'Date', 'week': 'Week', 'RUN ID': 'RUN ID',
        'run_label': 'RUN ID', 'stops': 'Stops', 'STOPS': 'Stops',
        'total_shots': 'Total Shots', 'Total Shots': 'Total Shots',
        'startup_shots': 'Start-up Shots', 'Start-up Shots': 'Start-up Shots',
        'mttr_min': 'MTTR (min)', 'MTTR (min)': 'MTTR (min)',
        'mtbf_min': 'MTBF (min)', 'MTBF (min)': 'MTBF (min)',
        'stability_index': 'Stability Index (%)', 'STABILITY %': 'Stability Index (%)',
        'approved_ct': 'Approved CT', 'APPROVED CT': 'Approved CT',
        'mode_ct': 'Mode CT', 'MODE CT': 'Mode CT'
    }

    cols_to_keep = [col for col in df.columns if col in rename_map]
    df_filtered = df[cols_to_keep]
    cols_to_rename = {k: v for k, v in rename_map.items() if k in df_filtered.columns}
    df_renamed = df_filtered.rename(columns=cols_to_rename)

    display_order = [
        'Hour', 'Date', 'Week', 'RUN ID', 'Approved CT', 'Mode CT',
        'Stops', 'Total Shots', 'Start-up Shots', 'Stability Index (%)', 'MTTR (min)', 'MTBF (min)'
    ]
    final_cols = [col for col in display_order if col in df_renamed.columns]
    for col in df_renamed.columns:
        if col not in final_cols:
            final_cols.append(col)

    return df_renamed[final_cols]


@st.cache_data
def load_all_data(files, date_format="ISO / Default (YYYY-MM-DD)", _cache_version=None):
    df_list = []
    for file in files:
        try:
            if file.name.endswith('.csv'):
                df = pd.read_csv(file, low_memory=False)
            else:
                df = pd.read_excel(file)

            col_map = {col.strip().upper(): col for col in df.columns}

            def get_col(*targets):
                for t in targets:
                    found = col_map.get(t.strip().upper())
                    if found is not None:
                        return found
                return None

            tool_id_col = get_col("EQUIPMENT_CODE", "TOOLING ID", "EQUIPMENT CODE", "TOOL_ID")
            if tool_id_col:
                df.rename(columns={tool_id_col: "tool_id"}, inplace=True)

            actual_ct_col = get_col("CT", "ACTUAL CT", "ACTUAL_CT", "CYCLE TIME")
            if actual_ct_col:
                df.rename(columns={actual_ct_col: "ACTUAL CT"}, inplace=True)

            approved_ct_col = get_col("APPROVED_CT", "APPROVED CT", "STANDARD CT", "STD CT")
            if approved_ct_col:
                df["approved_ct"] = pd.to_numeric(df[approved_ct_col], errors='coerce')
            else:
                df["approved_ct"] = np.nan

            if {"YEAR", "MONTH", "DAY", "TIME"}.issubset(set(col_map.keys())):
                datetime_str = (
                    df[col_map["YEAR"]].astype(str) + "-"
                    + df[col_map["MONTH"]].astype(str) + "-"
                    + df[col_map["DAY"]].astype(str) + " "
                    + df[col_map["TIME"]].astype(str)
                )
                df["shot_time"] = pd.to_datetime(datetime_str, errors="coerce")
            else:
                shot_time_col = get_col("LOCAL_SHOT_TIME", "SHOT TIME", "TIMESTAMP", "DATE", "TIME")
                if shot_time_col:
                    if date_format == "Global (DD/MM/YYYY)":
                        parsed = pd.to_datetime(df[shot_time_col], format="mixed", dayfirst=True, errors="coerce")
                    elif date_format == "US (MM/DD/YYYY)":
                        parsed = pd.to_datetime(df[shot_time_col], format="mixed", dayfirst=False, errors="coerce")
                    else:
                        # Ensures proper Year-Month-Day parsing
                        parsed = pd.to_datetime(df[shot_time_col], format="mixed", yearfirst=True, dayfirst=False, errors="coerce")
                        
                    df["shot_time"] = parsed

            session_col = get_col("COUNTER_CODE", "SESSION ID")
            if session_col and session_col != "SESSION ID":
                df.rename(columns={session_col: "SESSION ID"}, inplace=True)

            supplier_col = get_col("SUPPLIER_NAME", "SUPPLIER NAME", "SUPPLIER")
            if supplier_col and supplier_col != "supplier_name":
                df.rename(columns={supplier_col: "supplier_name"}, inplace=True)

            part_id_col = get_col("PART_ID", "PART")
            if part_id_col and part_id_col != "part_id":
                df.rename(columns={part_id_col: "part_id"}, inplace=True)

            part_name_col = get_col("PART_NAME", "PART NAME")
            if part_name_col and part_name_col != "part_name":
                df.rename(columns={part_name_col: "part_name"}, inplace=True)

            tooling_type_col = get_col("TOOLING_TYPE", "TOOLING TYPE")
            if tooling_type_col and tooling_type_col != "tooling_type":
                df.rename(columns={tooling_type_col: "tooling_type"}, inplace=True)

            project_col = get_col("PROJECT_ID", "PROJECT", "PROJECT NAME")
            if project_col and project_col != "project_id":
                df.rename(columns={project_col: "project_id"}, inplace=True)

            if "tool_id" in df.columns and "shot_time" in df.columns:
                df_list.append(df)
        except Exception as e:
            st.warning(f"Could not load file: {file.name}. Error: {e}")

    if not df_list:
        return pd.DataFrame()

    df_final = pd.concat(df_list, ignore_index=True)
    if 'tool_id' not in df_final.columns:
        df_final['tool_id'] = 'Unknown'
    df_final['tool_id'] = df_final['tool_id'].fillna('Unknown').astype(str)

    return df_final


# ==============================================================================
# --- 2. CORE CALCULATION ENGINE ---
# ==============================================================================

def _get_stable_mode(series: pd.Series) -> float:
    if series.empty:
        return 0.0
    rounded = series.round(2)
    modes = rounded.mode()
    return float(modes.iloc[0]) if not modes.empty else float(series.mean())


class RunRateCalculator:
    def __init__(self, df: pd.DataFrame, tolerance: float,
                 downtime_gap_tolerance: float, analysis_mode: str = 'aggregate',
                 run_interval_hours: float = 8, startup_shots_count: int = 0):
        self.df_raw = df.copy()
        self.tolerance = tolerance
        self.downtime_gap_tolerance = downtime_gap_tolerance
        self.analysis_mode = analysis_mode
        self.run_interval_hours = run_interval_hours
        self.startup_shots_count = startup_shots_count
        self.results = self._calculate_all_metrics()

    def _calculate_hourly_summary(self, df: pd.DataFrame) -> pd.DataFrame:
        if df.empty or 'stop_event' not in df.columns:
            return pd.DataFrame()

        df = df.copy()
        df['hour'] = df['shot_time'].dt.hour
        hourly_groups = df.groupby('hour')

        stops = hourly_groups['stop_event'].sum()
        hourly_total_downtime_sec = hourly_groups.apply(
            lambda x: x[x['stop_flag'] == 1]['adj_ct_sec'].sum()
        )
        uptime_min = df[df['stop_flag'] == 0].groupby('hour')['ACTUAL CT'].sum() / 60
        shots = hourly_groups.size().rename('total_shots')
        
        if 'startup_flag' in df.columns:
            startup_shots = hourly_groups['startup_flag'].sum().rename('startup_shots')
        else:
            startup_shots = pd.Series(0, index=hourly_groups.groups.keys(), name='startup_shots')

        hourly_summary = pd.DataFrame(index=range(24))
        hourly_summary['hour'] = hourly_summary.index
        hourly_summary = (hourly_summary
                          .join(stops.rename('stops'))
                          .join(shots)
                          .join(startup_shots)
                          .join(uptime_min.rename('uptime_min'))
                          .fillna(0)
                          .join(hourly_total_downtime_sec.rename('total_downtime_sec'))
                          .fillna(0))

        hourly_summary['mttr_min'] = ((hourly_summary['total_downtime_sec'] / 60)
                                      / hourly_summary['stops'].replace(0, np.nan))
        hourly_summary['mtbf_min'] = (hourly_summary['uptime_min']
                                      / hourly_summary['stops'].replace(0, np.nan))
        hourly_summary['mtbf_min'] = hourly_summary['mtbf_min'].fillna(hourly_summary['uptime_min'])

        effective_runtime_min = hourly_summary['uptime_min'] + (hourly_summary['total_downtime_sec'] / 60)
        hourly_summary['stability_index'] = np.where(
            effective_runtime_min > 0,
            (hourly_summary['uptime_min'] / effective_runtime_min) * 100,
            np.where(hourly_summary['stops'] == 0, 100.0, 0.0)
        )
        hourly_summary['stability_index'] = np.where(
            hourly_summary['total_shots'] == 0, np.nan, hourly_summary['stability_index']
        )

        if 'approved_ct' in df.columns:
            hourly_approved = hourly_groups['approved_ct'].apply(
                lambda x: x.mode().iloc[0] if not x.mode().empty else np.nan
            )
            hourly_summary = hourly_summary.join(hourly_approved.rename('approved_ct'))

        if 'mode_ct' in df.columns:
            hourly_mode = hourly_groups['ACTUAL CT'].apply(_get_stable_mode)
            hourly_summary = hourly_summary.join(hourly_mode.rename('mode_ct'))

        cols_to_fill = [col for col in hourly_summary.columns
                        if col not in ['stability_index', 'approved_ct', 'mode_ct']]
        hourly_summary[cols_to_fill] = hourly_summary[cols_to_fill].fillna(0)

        return hourly_summary

    def _calculate_all_metrics(self) -> dict:
        df = self.df_raw.copy()
        if df.empty or "shot_time" not in df.columns:
            return {}

        if 'ACTUAL CT' not in df.columns:
            df['ACTUAL CT'] = np.nan
        df['ACTUAL CT'] = pd.to_numeric(df['ACTUAL CT'], errors='coerce')
        df = (df.dropna(subset=['shot_time', 'ACTUAL CT'])
                .sort_values(['tool_id', 'shot_time', 'ACTUAL CT'])
                .reset_index(drop=True))
        if df.empty:
            return {}

        df['time_diff_sec'] = (df.groupby('tool_id')['shot_time']
                               .diff().dt.total_seconds().fillna(0))
        mask_first_shot = df['tool_id'] != df['tool_id'].shift(1)
        df.loc[mask_first_shot, 'time_diff_sec'] = df.loc[mask_first_shot, 'ACTUAL CT']

        is_new_run = df['time_diff_sec'] > (self.run_interval_hours * 3600)
        df['run_id'] = (is_new_run | mask_first_shot).cumsum()

        # START-UP SHOT LOGIC
        df['shot_seq_in_run'] = df.groupby('run_id').cumcount() + 1
        df['startup_flag'] = np.where(df['shot_seq_in_run'] <= self.startup_shots_count, 1, 0)
        df['startup_event'] = (df['startup_flag'] == 1)

        run_modes = (df[df['ACTUAL CT'] < 1000]
                     .groupby('run_id')['ACTUAL CT']
                     .apply(_get_stable_mode))
        df['mode_ct'] = df['run_id'].map(run_modes)
        df['mode_ct'] = df['mode_ct'].fillna(df['ACTUAL CT'].median())

        df['lower_limit'] = df['mode_ct'] * (1 - self.tolerance)
        df['upper_limit'] = df['mode_ct'] * (1 + self.tolerance)

        modes_unique = df['mode_ct'].dropna().unique()
        mode_ct_display = (float(modes_unique[0])
                           if len(modes_unique) == 1
                           else "Varies by Run")

        approved_ct_display = np.nan
        if 'approved_ct' in df.columns:
            valid_app = df['approved_ct'].dropna()
            if not valid_app.empty:
                approved_ct_display = (valid_app.mode().iloc[0]
                                       if not valid_app.mode().empty
                                       else valid_app.mean())

        df['next_shot_time_diff'] = df.groupby('tool_id')['time_diff_sec'].shift(-1).fillna(0)
        is_time_gap = df['next_shot_time_diff'] > (df['ACTUAL CT'] + self.downtime_gap_tolerance)
        is_abnormal = (df['ACTUAL CT'] < df['lower_limit']) | (df['ACTUAL CT'] > df['upper_limit'])
        is_hard_stop = df['ACTUAL CT'] >= 999.9

        df['stop_flag'] = np.where(is_time_gap | is_abnormal | is_hard_stop, 1, 0)
        df['prev_stop_flag'] = df.groupby('tool_id')['stop_flag'].shift(1, fill_value=0)
        df['stop_event'] = (df['stop_flag'] == 1) & (df['prev_stop_flag'] == 0)

        df['adj_ct_sec'] = df['ACTUAL CT']
        df.loc[is_time_gap, 'adj_ct_sec'] = df['next_shot_time_diff']

        run_durations_sec = []
        for _, run_df in df.groupby('run_id'):
            if not run_df.empty:
                start = run_df['shot_time'].min()
                end = run_df['shot_time'].max()
                last_ct = run_df.iloc[-1]['ACTUAL CT']
                run_durations_sec.append((end - start).total_seconds() + last_ct)

        total_runtime_sec = sum(run_durations_sec)
        prod_df = df[df['stop_flag'] == 0]
        production_time_sec = prod_df['ACTUAL CT'].sum()
        downtime_sec = max(0, total_runtime_sec - production_time_sec)

        total_shots = len(df)
        stop_events = df['stop_event'].sum()
        
        # Recalculated Bucket Logic 
        startup_shots = df['startup_flag'].sum()
        startup_within_mode = ((df['startup_flag'] == 1) & (df['ACTUAL CT'] >= df['lower_limit']) & (df['ACTUAL CT'] <= df['upper_limit'])).sum()
        startup_outside_mode = startup_shots - startup_within_mode
        
        normal_shots = len(df[(df['stop_flag'] == 0) & (df['startup_flag'] == 0)])
        stop_shots = len(df[(df['stop_flag'] == 1) & (df['startup_flag'] == 0)])
        non_stop_shots = len(prod_df)

        mttr_min = (downtime_sec / 60 / stop_events) if stop_events > 0 else 0
        mtbf_min = ((production_time_sec / 60 / stop_events)
                    if stop_events > 0
                    else (production_time_sec / 60))
        stability_index = ((production_time_sec / total_runtime_sec * 100)
                           if total_runtime_sec > 0 else 100.0)
        efficiency = (non_stop_shots / total_shots) if total_shots > 0 else 0
        avg_cycle_time_sec = production_time_sec / non_stop_shots if non_stop_shots > 0 else 0

        first_stop_event_index = df[df['stop_event'] == True].index.min()
        if pd.isna(first_stop_event_index):
            time_to_first_dt_sec = production_time_sec
        elif first_stop_event_index == 0:
            time_to_first_dt_sec = 0
        else:
            time_to_first_dt_sec = df.loc[:first_stop_event_index - 1, 'adj_ct_sec'].sum()

        df["run_group"] = df["stop_event"].cumsum()
        df_for_runs = df[df['adj_ct_sec'] <= 28800].copy()
        run_durations = (df_for_runs[df_for_runs["stop_flag"] == 0]
                         .groupby("run_group")["ACTUAL CT"]
                         .sum().div(60)
                         .reset_index(name="duration_min"))

        max_minutes = min(run_durations["duration_min"].max(), 240) if not run_durations.empty else 0
        upper_bound = int(np.ceil(max_minutes / 20.0) * 20)
        edges = list(range(0, upper_bound + 20, 20)) if upper_bound > 0 else [0, 20]
        labels = [f"{edges[i]} to <{edges[i+1]}" for i in range(len(edges) - 1)]
        if labels:
            labels[-1] = f"{edges[-2]}+"
        if edges and len(edges) > 1:
            edges[-1] = np.inf

        if not run_durations.empty:
            run_durations["time_bucket"] = pd.cut(
                run_durations["duration_min"], bins=edges,
                labels=labels, right=False, include_lowest=True
            )

        reds = px.colors.sequential.Reds[3:7][::-1]
        blues = px.colors.sequential.Blues[3:8]
        greens = px.colors.sequential.Greens[3:8]
        red_labels, blue_labels, green_labels = [], [], []
        for label in labels:
            try:
                lower_bound = int(label.split(' ')[0].replace('+', ''))
                if lower_bound < 60:
                    red_labels.append(label)
                elif 60 <= lower_bound < 160:
                    blue_labels.append(label)
                else:
                    green_labels.append(label)
            except (ValueError, IndexError):
                continue

        bucket_color_map = {}
        for i, label in enumerate(red_labels):
            bucket_color_map[label] = reds[i % len(reds)]
        for i, label in enumerate(blue_labels):
            bucket_color_map[label] = blues[i % len(blues)]
        for i, label in enumerate(green_labels):
            bucket_color_map[label] = greens[i % len(greens)]

        hourly_summary = self._calculate_hourly_summary(df)

        return {
            "processed_df": df,
            "mode_ct": mode_ct_display,
            "total_shots": total_shots,
            "efficiency": efficiency,
            "stop_events": stop_events,
            "normal_shots": normal_shots,
            "startup_shots": startup_shots,
            "startup_within_mode": startup_within_mode,
            "startup_outside_mode": startup_outside_mode,
            "stopped_shots": stop_shots,
            "non_stop_shots": non_stop_shots,
            "mttr_min": mttr_min,
            "mtbf_min": mtbf_min,
            "stability_index": stability_index,
            "run_durations": run_durations,
            "bucket_labels": labels,
            "bucket_color_map": bucket_color_map,
            "hourly_summary": hourly_summary,
            "total_runtime_sec": total_runtime_sec,
            "production_time_sec": production_time_sec,
            "downtime_sec": downtime_sec,
            "avg_cycle_time_sec": avg_cycle_time_sec,
            "time_to_first_dt_min": time_to_first_dt_sec / 60,
            "production_run_sec": total_runtime_sec,
            "tot_down_time_sec": downtime_sec,
            "approved_ct": approved_ct_display,
            "lower_limit": (df['lower_limit'].min() if not df.empty else 0),
            "upper_limit": (df['upper_limit'].max() if not df.empty else 0),
            "min_lower_limit": df['lower_limit'].min() if not df.empty else 0,
            "max_lower_limit": df['lower_limit'].max() if not df.empty else 0,
            "min_upper_limit": df['upper_limit'].min() if not df.empty else 0,
            "max_upper_limit": df['upper_limit'].max() if not df.empty else 0,
            "min_mode_ct": (df['mode_ct'].min() if not df.empty and pd.notna(df['mode_ct'].min()) else 0),
            "max_mode_ct": (df['mode_ct'].max() if not df.empty and pd.notna(df['mode_ct'].max()) else 0),
        }


def _run_metrics_from_processed(df_slice: pd.DataFrame) -> dict:
    if df_slice.empty:
        return {}
    start = df_slice['shot_time'].min()
    end = df_slice['shot_time'].max()
    last_ct = df_slice.iloc[-1]['ACTUAL CT']
    duration = (end - start).total_seconds() + last_ct

    prod_df = df_slice[df_slice['stop_flag'] == 0]
    prod_sec = prod_df['ACTUAL CT'].sum()
    down_sec = max(0, duration - prod_sec)
    tot_stops = df_slice['stop_event'].sum()
    tot_shots = len(df_slice)

    if 'startup_flag' in df_slice.columns:
        startup_shots = int(df_slice['startup_flag'].sum())
        startup_within = int(((df_slice['startup_flag'] == 1) & 
                              (df_slice['ACTUAL CT'] >= df_slice['lower_limit']) & 
                              (df_slice['ACTUAL CT'] <= df_slice['upper_limit'])).sum())
        startup_outside = startup_shots - startup_within
        normal_shots = int(len(df_slice[(df_slice['stop_flag'] == 0) & (df_slice['startup_flag'] == 0)]))
        stopped_shots = int(len(df_slice[(df_slice['stop_flag'] == 1) & (df_slice['startup_flag'] == 0)]))
    else:
        startup_shots = 0
        startup_within = 0
        startup_outside = 0
        normal_shots = len(prod_df)
        stopped_shots = tot_shots - normal_shots

    if 'mode_ct' in df_slice.columns and not df_slice['mode_ct'].dropna().empty:
        mode_ct = float(df_slice['mode_ct'].iloc[0])
    else:
        mode_ct = _get_stable_mode(df_slice['ACTUAL CT']) if tot_shots > 0 else np.nan
        
    approved_ct = (df_slice['approved_ct'].mode().iloc[0]
                   if 'approved_ct' in df_slice.columns
                   and not df_slice['approved_ct'].dropna().empty
                   else np.nan)

    return {
        'start': start, 'end': end,
        'duration': duration,
        'prod_sec': prod_sec,
        'down_sec': down_sec,
        'tot_stops': tot_stops,
        'tot_shots': tot_shots,
        'normal_shots': normal_shots,
        'startup_shots': startup_shots,
        'startup_within': startup_within,
        'startup_outside': startup_outside,
        'stopped_shots': stopped_shots,
        'non_stop_shots': len(prod_df),
        'mode_ct': mode_ct,
        'approved_ct': approved_ct,
        'stability_index': (prod_sec / duration * 100) if duration > 0 else 100.0,
        'mttr_min': (down_sec / 60 / tot_stops) if tot_stops > 0 else np.nan,
        'mtbf_min': ((prod_sec / 60 / tot_stops) if tot_stops > 0 else (prod_sec / 60)),
    }


def calculate_daily_summaries_for_week(df_week, tolerance, downtime_gap_tolerance,
                                       analysis_mode, run_interval_hours=8, startup_shots_count=0):
    calc_global = RunRateCalculator(df_week, tolerance, downtime_gap_tolerance,
                                    analysis_mode, run_interval_hours, startup_shots_count)
    df_proc = calc_global.results.get('processed_df', df_week)
    if df_proc.empty:
        return pd.DataFrame()

    df_proc = df_proc.copy()
    df_proc['date_lbl'] = df_proc['shot_time'].dt.date

    daily_results_list = []
    for date_val, df_day in df_proc.groupby('date_lbl'):
        if df_day.empty:
            continue
        m = _run_metrics_from_processed(df_day)
        daily_results_list.append({
            'date': date_val,
            'stability_index': m['stability_index'],
            'mttr_min': m['mttr_min'],
            'mtbf_min': m['mtbf_min'],
            'stops': m['tot_stops'],
            'total_shots': m['tot_shots'],
            'total_downtime_sec': m['down_sec'],
            'uptime_min': m['prod_sec'] / 60,
            'mode_ct': m['mode_ct'],
            'approved_ct': m['approved_ct'],
        })
    return pd.DataFrame(daily_results_list)


def calculate_weekly_summaries_for_month(df_month, tolerance, downtime_gap_tolerance,
                                         analysis_mode, run_interval_hours=8, startup_shots_count=0):
    calc_global = RunRateCalculator(df_month, tolerance, downtime_gap_tolerance,
                                    analysis_mode, run_interval_hours, startup_shots_count)
    df_proc = calc_global.results.get('processed_df', df_month)
    if df_proc.empty:
        return pd.DataFrame()

    df_proc = df_proc.copy()
    iso = df_proc['shot_time'].dt.isocalendar()
    df_proc['week_lbl'] = iso['year'].astype(str) + '-W' + iso['week'].astype(str).str.zfill(2)

    weekly_results_list = []
    for week, df_week in df_proc.groupby('week_lbl'):
        if df_week.empty:
            continue
        m = _run_metrics_from_processed(df_week)
        weekly_results_list.append({
            'week': week,
            'stability_index': m['stability_index'],
            'mttr_min': m['mttr_min'],
            'mtbf_min': m['mtbf_min'],
            'stops': m['tot_stops'],
            'total_shots': m['tot_shots'],
            'total_downtime_sec': m['down_sec'],
            'uptime_min': m['prod_sec'] / 60,
            'mode_ct': m['mode_ct'],
            'approved_ct': m['approved_ct'],
        })
    return pd.DataFrame(weekly_results_list)


def build_display_results(df: pd.DataFrame, run_interval_hours: float = 8) -> dict:
    if df.empty:
        return {}

    col = 'adj_ct_sec' if 'adj_ct_sec' in df.columns else 'ACTUAL CT'
    df_for_runs = df[df[col] <= 28800].copy()
    run_durations = (
        df_for_runs[df_for_runs["stop_flag"] == 0]
        .groupby("run_group")["ACTUAL CT"]
        .sum().div(60)
        .reset_index(name="duration_min")
    )

    max_minutes = min(run_durations["duration_min"].max(), 240) if not run_durations.empty else 0
    upper_bound = int(np.ceil(max_minutes / 20.0) * 20)
    edges = list(range(0, upper_bound + 20, 20)) if upper_bound > 0 else [0, 20]
    labels = [f"{edges[i]} to <{edges[i+1]}" for i in range(len(edges) - 1)]
    if labels:
        labels[-1] = f"{edges[-2]}+"
    if edges and len(edges) > 1:
        edges[-1] = np.inf

    if not run_durations.empty:
        run_durations["time_bucket"] = pd.cut(
            run_durations["duration_min"], bins=edges,
            labels=labels, right=False, include_lowest=True
        )

    reds = px.colors.sequential.Reds[3:7][::-1]
    blues = px.colors.sequential.Blues[3:8]
    greens = px.colors.sequential.Greens[3:8]
    red_labels, blue_labels, green_labels = [], [], []
    for label in labels:
        try:
            lb = int(label.split(' ')[0].replace('+', ''))
            if lb < 60:
                red_labels.append(label)
            elif lb < 160:
                blue_labels.append(label)
            else:
                green_labels.append(label)
        except (ValueError, IndexError):
            continue

    bucket_color_map = {}
    for i, label in enumerate(red_labels):
        bucket_color_map[label] = reds[i % len(reds)]
    for i, label in enumerate(blue_labels):
        bucket_color_map[label] = blues[i % len(blues)]
    for i, label in enumerate(green_labels):
        bucket_color_map[label] = greens[i % len(greens)]

    modes_unique = df['mode_ct'].dropna().unique() if 'mode_ct' in df.columns else np.array([])
    if len(modes_unique) == 1:
        mode_ct_display = float(modes_unique[0])
    elif len(modes_unique) > 1:
        mode_ct_display = "Varies by Run"
    else:
        mode_ct_display = 0

    lower_limit = df['lower_limit'].min() if 'lower_limit' in df.columns else 0
    upper_limit = df['upper_limit'].max() if 'upper_limit' in df.columns else 0

    temp = object.__new__(RunRateCalculator)
    hourly_summary = temp._calculate_hourly_summary(df)

    return {
        "processed_df": df,
        "run_durations": run_durations,
        "bucket_labels": labels,
        "bucket_color_map": bucket_color_map,
        "mode_ct": mode_ct_display,
        "lower_limit": lower_limit,
        "upper_limit": upper_limit,
        "hourly_summary": hourly_summary,
    }


def calculate_run_summaries(df_period, tolerance, downtime_gap_tolerance,
                            run_interval_hours=8, pre_processed=False, startup_shots_count=0):
    if pre_processed and 'stop_flag' in df_period.columns:
        df_proc = df_period
    else:
        calc_base = RunRateCalculator(df_period, tolerance, downtime_gap_tolerance,
                                      'aggregate', run_interval_hours, startup_shots_count)
        df_proc = calc_base.results.get('processed_df', df_period)

    run_summary_list = []
    for seq, (run_id_val, df_run) in enumerate(df_proc.groupby('run_id'), start=1):
        if df_run.empty:
            continue
        m = _run_metrics_from_processed(df_run)

        lower_limit = df_run['lower_limit'].iloc[0] if 'lower_limit' in df_run.columns else 0
        upper_limit = df_run['upper_limit'].iloc[0] if 'upper_limit' in df_run.columns else 0

        run_summary_list.append({
            'run_id': run_id_val,
            'run_label': f"Run {seq:03d}",
            'start_time': m['start'],
            'end_time': m['end'],
            'total_shots': m['tot_shots'],
            'normal_shots': m['normal_shots'],
            'startup_shots': m.get('startup_shots', 0),
            'startup_within': m.get('startup_within', 0),
            'startup_outside': m.get('startup_outside', 0),
            'stopped_shots': m.get('stopped_shots', m['tot_shots'] - m['normal_shots']),
            'non_stop_shots': m.get('non_stop_shots', m['normal_shots']),
            'mode_ct': m['mode_ct'],
            'lower_limit': lower_limit,
            'upper_limit': upper_limit,
            'total_runtime_sec': m['duration'],
            'production_time_sec': m['prod_sec'],
            'downtime_sec': m['down_sec'],
            'mttr_min': (m['down_sec'] / 60 / m['tot_stops']) if m['tot_stops'] > 0 else 0,
            'mtbf_min': ((m['prod_sec'] / 60 / m['tot_stops']) if m['tot_stops'] > 0 else (m['prod_sec'] / 60)),
            'stability_index': m['stability_index'],
            'stops': m['tot_stops'],
            'approved_ct': m['approved_ct'],
        })

    if not run_summary_list:
        return pd.DataFrame()
    return (pd.DataFrame(run_summary_list)
              .sort_values('start_time')
              .reset_index(drop=True))


# ==============================================================================
# --- 3. PLOTTING FUNCTIONS ---
# ==============================================================================

def create_gauge(value, title, steps=None):
    color = "#3498DB"
    if steps:
        if value <= 50:
            color = PASTEL_COLORS['red']
        elif value <= 70:
            color = PASTEL_COLORS['orange']
        else:
            color = PASTEL_COLORS['green']

    plot_value = max(0, min(value, 100))
    remainder = 100 - plot_value

    fig = go.Figure(data=[go.Pie(
        values=[plot_value, remainder], hole=0.75, sort=False, direction='clockwise',
        textinfo='none', marker=dict(colors=[color, '#e6e6e6']), hoverinfo='none'
    )])
    fig.add_annotation(
        text=f"{value:.1f}%", x=0.5, y=0.5,
        font=dict(size=42, weight='bold', color=color, family="Arial"),
        showarrow=False
    )
    fig.update_layout(
        title=dict(text=title, x=0.5, xanchor='center', y=0.95, font=dict(size=16)),
        margin=dict(l=20, r=20, t=40, b=20), height=250, showlegend=False,
        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)'
    )
    return fig


def plot_shot_bar_chart(df, lower_limit, upper_limit, mode_ct,
                        time_agg='hourly', show_approved_ct=False):
    if df.empty:
        st.info("No shot data to display for this period.")
        return
    df = df.copy()
    
    df['color'] = np.where(df['startup_flag'] == 1, 'purple', 
                  np.where(df['stop_flag'] == 1, PASTEL_COLORS['red'], '#3498DB'))

    downtime_gap_indices = df[df['adj_ct_sec'] != df['ACTUAL CT']].index
    valid_downtime_gap_indices = downtime_gap_indices[downtime_gap_indices > 0]
    normal_shot_indices = df.index.difference(valid_downtime_gap_indices)

    if not normal_shot_indices.empty:
        shot_index_in_second = df.loc[normal_shot_indices].groupby('shot_time').cumcount()
        time_offset = pd.to_timedelta(shot_index_in_second * 0.2, unit='s')
        df.loc[normal_shot_indices, 'plot_time'] = (
            df.loc[normal_shot_indices, 'shot_time'] + time_offset
        )

    if not valid_downtime_gap_indices.empty:
        prev_shot_timestamps = df['shot_time'].shift(1).loc[valid_downtime_gap_indices]
        df.loc[valid_downtime_gap_indices, 'plot_time'] = prev_shot_timestamps

    _first_idx = df.index[0]
    if pd.isna(df.loc[_first_idx, 'plot_time']) if 'plot_time' in df.columns else True:
        df.loc[_first_idx, 'plot_time'] = df.loc[_first_idx, 'shot_time']

    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['plot_time'], y=df['adj_ct_sec'],
        marker_color=df['color'], name='Cycle Time', showlegend=False
    ))
    fig.add_trace(go.Bar(x=[None], y=[None], name="Normal Shot",
                         marker_color='#3498DB', showlegend=True))
    fig.add_trace(go.Bar(x=[None], y=[None], name="Stopped Shot",
                         marker_color=PASTEL_COLORS['red'], showlegend=True))
    fig.add_trace(go.Bar(x=[None], y=[None], name="Start-up Shot",
                         marker_color='purple', showlegend=True))
    fig.add_trace(go.Scatter(
        x=[None], y=[None], mode='lines', line=dict(width=0),
        fill='tozeroy', fillcolor='rgba(119, 221, 119, 0.3)',
        name='Tolerance Band', showlegend=True
    ))
    fig.add_trace(go.Scatter(
        x=[None], y=[None], mode='lines', name='New Run Start',
        line=dict(color='purple', dash='dash', width=2), showlegend=True
    ))

    if show_approved_ct and 'approved_ct' in df.columns:
        fig.add_trace(go.Scatter(
            x=df['plot_time'], y=df['approved_ct'], mode='lines',
            name='Approved CT', line=dict(color='#00FF00', width=2, dash='dash')
        ))

    if 'lower_limit' in df.columns and 'run_id' in df.columns:
        for run_id_val, group in df.groupby('run_id'):
            if not group.empty:
                fig.add_shape(
                    type="rect", xref="x", yref="y",
                    x0=group['shot_time'].min(), y0=group['lower_limit'].iloc[0],
                    x1=group['shot_time'].max(), y1=group['upper_limit'].iloc[0],
                    fillcolor=PASTEL_COLORS['green'], opacity=0.3,
                    layer="below", line_width=0
                )
    else:
        if not df.empty:
            fig.add_shape(
                type="rect", xref="x", yref="y",
                x0=df['shot_time'].min(), y0=lower_limit,
                x1=df['shot_time'].max(), y1=upper_limit,
                fillcolor=PASTEL_COLORS['green'], opacity=0.3,
                layer="below", line_width=0
            )

    if 'run_id' in df.columns:
        run_starts = df.groupby('run_id')['shot_time'].min().sort_values()
        for start_time in run_starts.iloc[1:]:
            fig.add_vline(x=start_time, line_width=2, line_dash="dash", line_color="purple")

    y_axis_cap_val = (mode_ct if isinstance(mode_ct, (int, float))
                      else df['mode_ct'].mean() if 'mode_ct' in df else 50)
    y_axis_cap = min(max(y_axis_cap_val * 2, 50), 500)

    fig.update_layout(
        title="Run Rate Cycle Time", xaxis_title="Date / Time",
        yaxis_title="Cycle Time (sec)", yaxis=dict(range=[0, y_axis_cap]),
        bargap=0.05, xaxis=dict(showgrid=True), showlegend=True,
        legend=dict(title="Legend", orientation="h", yanchor="bottom",
                    y=1.02, xanchor="right", x=1)
    )
    st.plotly_chart(fig, width='stretch')


def plot_trend_chart(df, x_col, y_col, title, x_title, y_title,
                     y_range=None, is_stability=False):
    if y_range is None:
        y_range = [0, 101]
    fig = go.Figure()
    marker_config = {}

    if y_col not in df.columns:
        return
    plot_df = df.dropna(subset=[y_col])
    if plot_df.empty:
        return

    if is_stability:
        marker_config['color'] = [
            PASTEL_COLORS['red'] if v <= 50
            else PASTEL_COLORS['orange'] if v <= 70
            else PASTEL_COLORS['green']
            for v in plot_df[y_col]
        ]
        marker_config['size'] = 10

    fig.add_trace(go.Scatter(
        x=plot_df[x_col], y=plot_df[y_col], mode="lines+markers",
        name=y_title,
        line=dict(color="black" if is_stability else "royalblue", width=2),
        marker=marker_config
    ))
    if is_stability:
        for y0, y1, c in [(0, 50, PASTEL_COLORS['red']),
                          (50, 70, PASTEL_COLORS['orange']),
                          (70, 100, PASTEL_COLORS['green'])]:
            fig.add_shape(
                type="rect", xref="paper", x0=0, x1=1, y0=y0, y1=y1,
                fillcolor=c, opacity=0.2, line_width=0, layer="below"
            )

    fig.update_layout(
        title=title,
        yaxis=dict(title=y_title, range=y_range),
        xaxis_title=x_title.title(),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
    )
    st.plotly_chart(fig, width='stretch')


def plot_mttr_mtbf_chart(df, x_col, mttr_col, mtbf_col, shots_col, title):
    if df is None or df.empty or df[shots_col].sum() == 0:
        return
    required_cols = [x_col, mttr_col, mtbf_col, shots_col]
    if not all(col in df.columns for col in required_cols):
        return

    mttr = df[mttr_col]
    mtbf = df[mtbf_col]
    shots = df[shots_col]
    x_axis = df[x_col]

    max_mttr = np.nanmax(mttr[np.isfinite(mttr)]) if any(np.isfinite(mttr)) else 0
    max_mtbf = np.nanmax(mtbf[np.isfinite(mtbf)]) if any(np.isfinite(mtbf)) else 0
    y_range_mttr = [0, max_mttr * 1.15 if max_mttr > 0 else 10]
    y_range_mtbf = [0, max_mtbf * 1.15 if max_mtbf > 0 else 10]

    shots_min, shots_max = shots.min(), shots.max()
    if (shots_max - shots_min) == 0:
        scaled_shots = pd.Series(
            [y_range_mtbf[1] / 2 if y_range_mtbf[1] > 0 else 0.5] * len(shots),
            index=shots.index
        )
    else:
        scaled_shots = ((shots - shots_min) / (shots_max - shots_min)
                        * (y_range_mtbf[1] * 0.9))

    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(go.Scatter(
        x=x_axis, y=mttr, name='MTTR (min)', mode='lines+markers',
        line=dict(color='red', width=4)
    ), secondary_y=False)
    fig.add_trace(go.Scatter(
        x=x_axis, y=mtbf, name='MTBF (min)', mode='lines+markers',
        line=dict(color='green', width=4)
    ), secondary_y=True)
    fig.add_trace(go.Scatter(
        x=x_axis, y=scaled_shots, name='Total Shots',
        mode='lines+markers+text', text=shots, textposition='top center',
        textfont=dict(color='blue'), line=dict(color='blue', dash='dot')
    ), secondary_y=True)

    fig.update_layout(
        title_text=title,
        yaxis_title="MTTR (min)", yaxis2_title="MTBF (min)",
        xaxis_title=x_col.replace("_", " ").title(),
        yaxis=dict(range=y_range_mttr), yaxis2=dict(range=y_range_mtbf),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
    )
    if x_col == 'hour':
        fig.update_layout(xaxis_title="Hour")
    st.plotly_chart(fig, width='stretch')


# ==============================================================================
# --- 4. TEXT ANALYSIS ENGINE ---
# ==============================================================================

def generate_detailed_analysis(analysis_df, overall_stability, overall_mttr,
                                overall_mtbf, analysis_level):
    if analysis_df is None or analysis_df.empty:
        return {"error": "Not enough data to generate a trend analysis."}

    stability_class = (
        "good (above 70%)" if overall_stability > 70
        else "needs improvement (50-70%)" if overall_stability > 50
        else "poor (below 50%)"
    )
    overall_summary = (
        f"The overall stability for this period is "
        f"<strong>{overall_stability:.1f}%</strong>, which is considered "
        f"<strong>{stability_class}</strong>."
    )

    predictive_insight = "Insufficient data points (less than 2 periods) to establish a predictive trend."
    analysis_df_clean = analysis_df.dropna(subset=['stability'])
    if len(analysis_df_clean) > 1:
        volatility_std = analysis_df_clean['stability'].std()
        volatility_level = (
            "highly volatile" if volatility_std > 15
            else "moderately volatile" if volatility_std > 5
            else "relatively stable"
        )
        half_point = len(analysis_df_clean) // 2
        first_half_mean = analysis_df_clean['stability'].iloc[:half_point].mean()
        second_half_mean = analysis_df_clean['stability'].iloc[half_point:].mean()

        if second_half_mean > first_half_mean * 1.05:
            trend_direction = "improving"
        elif second_half_mean < first_half_mean * 0.95:
            trend_direction = "declining"
        else:
            trend_direction = "stable"

        if trend_direction == "stable":
            predictive_insight = (
                f"Performance has been <strong>{volatility_level}</strong> "
                f"with no clear long-term upward or downward trend."
            )
        else:
            predictive_insight = (
                f"Performance shows a <strong>{trend_direction} trend</strong>, "
                f"although this has been <strong>{volatility_level}</strong>."
            )

    best_worst_analysis = ""
    if not analysis_df_clean.empty:
        best_performer = analysis_df_clean.loc[analysis_df_clean['stability'].idxmax()]
        worst_performer = analysis_df_clean.loc[analysis_df_clean['stability'].idxmin()]

        def format_period(period_value, level):
            if isinstance(period_value, (pd.Timestamp, pd.Period, pd.Timedelta, date, datetime)):
                return pd.to_datetime(period_value).strftime('%A, %b %d')
            if level == "Monthly":
                return f"Week {period_value}"
            if "Daily" in level:
                return f"{period_value}:00"
            return str(period_value)

        best_worst_analysis = (
            f"The best performance was during "
            f"<strong>{format_period(best_performer['period'], analysis_level)}</strong> "
            f"(Stability: {best_performer['stability']:.1f}%), "
            f"while the worst was during "
            f"<strong>{format_period(worst_performer['period'], analysis_level)}</strong> "
            f"(Stability: {worst_performer['stability']:.1f}%). "
            f"The key difference was the impact of stoppages: the worst period had "
            f"{int(worst_performer['stops'])} stops with an average duration of "
            f"{worst_performer.get('mttr', 0):.1f} min, compared to "
            f"{int(best_performer['stops'])} stops with an average duration of "
            f"{best_performer.get('mttr', 0):.1f} min during the best period."
        )

    pattern_insight = ""
    if not analysis_df_clean.empty and analysis_df_clean['stops'].sum() > 0:
        if "Daily" in analysis_level:
            peak_stop_hour = analysis_df_clean.loc[analysis_df_clean['stops'].idxmax()]
            try:
                period_val = int(float(peak_stop_hour['period']))
            except (ValueError, TypeError):
                period_val = str(peak_stop_hour['period'])
            try:
                stops_val = int(float(peak_stop_hour['stops']))
            except (ValueError, TypeError):
                stops_val = 0
            pattern_insight = (
                f"A notable pattern is the concentration of stop events around "
                f"<strong>{period_val}:00</strong>, which saw the highest number "
                f"of interruptions ({stops_val} stops)."
            )
        else:
            mean_stability = analysis_df_clean['stability'].mean()
            std_stability = analysis_df_clean['stability'].std()
            outlier_threshold = mean_stability - (1.5 * std_stability)
            outliers = analysis_df_clean[analysis_df_clean['stability'] < outlier_threshold]
            if not outliers.empty:
                worst_outlier = outliers.loc[outliers['stability'].idxmin()]
                pattern_insight = (
                    f"A key area of concern is "
                    f"<strong>{format_period(worst_outlier['period'], analysis_level)}</strong>, "
                    f"which performed significantly below average and disproportionately "
                    f"affected the overall stability."
                )

    if overall_stability >= 95:
        recommendation = "Overall performance is excellent. Continue monitoring for any emerging negative trends in either MTBF or MTTR."
    elif overall_stability > 70:
        if overall_mtbf > 0 and overall_mttr > 0 and overall_mtbf < (overall_mttr * 5):
            recommendation = f"Performance is good, but could be improved by focusing on <strong>Mean Time Between Failures (MTBF)</strong>. With an MTBF of <strong>{overall_mtbf:.1f} minutes</strong>, investigating the root causes of the more frequent, smaller stops could yield significant gains."
        else:
            recommendation = f"Performance is good, but could be improved by focusing on <strong>Mean Time To Repair (MTTR)</strong>. With an MTTR of <strong>{overall_mttr:.1f} minutes</strong>, streamlining the repair process for the infrequent but longer stops could yield significant gains."
    else:
        if overall_mtbf > 0 and overall_mttr > 0 and overall_mtbf < overall_mttr:
            recommendation = f"Stability is poor and requires attention. The primary driver is a low <strong>Mean Time Between Failures (MTBF)</strong> of <strong>{overall_mtbf:.1f} minutes</strong>. The top priority should be investigating the root cause of frequent machine stoppages."
        else:
            recommendation = f"Stability is poor and requires attention. The primary driver is a high <strong>Mean Time To Repair (MTTR)</strong> of <strong>{overall_mttr:.1f} minutes</strong>. The top priority should be investigating why stops take a long time to resolve."

    return {
        "overall": overall_summary,
        "predictive": predictive_insight,
        "best_worst": best_worst_analysis,
        "patterns": pattern_insight,
        "recommendation": recommendation,
    }


def generate_bucket_analysis(complete_runs, bucket_labels):
    if complete_runs.empty or 'duration_min' not in complete_runs.columns:
        return "No completed runs to analyze for long-run trends."

    total_completed_runs = len(complete_runs)
    try:
        long_run_buckets = [label for label in bucket_labels if int(label.split(' ')[0].replace('+', '')) >= 60]
    except (ValueError, IndexError):
        long_run_buckets = []

    num_long_runs = (complete_runs[complete_runs['time_bucket'].isin(long_run_buckets)].shape[0] if long_run_buckets else 0)
    percent_long_runs = (num_long_runs / total_completed_runs * 100) if total_completed_runs > 0 else 0
    longest_run_formatted = format_minutes_to_dhm(complete_runs['duration_min'].max())

    analysis_text = (
        f"Out of <strong>{total_completed_runs}</strong> completed runs, "
        f"<strong>{num_long_runs}</strong> ({percent_long_runs:.1f}%) qualified as long runs "
        f"(lasting over 60 minutes). The single longest stable run during this period lasted "
        f"for <strong>{longest_run_formatted}</strong>."
    )

    if total_completed_runs > 0:
        if percent_long_runs < 20:
            analysis_text += " This suggests that most stoppages occur after relatively short periods of operation."
        elif percent_long_runs > 50:
            analysis_text += " This indicates a strong capability for sustained stable operation."
        else:
            analysis_text += " This shows a mixed performance, with a reasonable number of long runs but also frequent shorter ones."
    return analysis_text


def generate_mttr_mtbf_analysis(analysis_df, analysis_level):
    analysis_df_clean = analysis_df.dropna(subset=['stops', 'stability', 'mttr'])
    if (analysis_df_clean.empty or analysis_df_clean['stops'].sum() == 0 or len(analysis_df_clean) < 2):
        return "Not enough stoppage data to generate a detailed correlation analysis."

    stops_stability_corr = analysis_df_clean['stops'].corr(analysis_df_clean['stability'])
    mttr_stability_corr = analysis_df_clean['mttr'].corr(analysis_df_clean['stability'])

    corr_insight = ""
    primary_driver_is_frequency = False
    primary_driver_is_duration = False

    if not pd.isna(stops_stability_corr) and not pd.isna(mttr_stability_corr):
        if abs(stops_stability_corr) > abs(mttr_stability_corr) * 1.5:
            primary_driver = "the **frequency of stops**"
            primary_driver_is_frequency = True
        elif abs(mttr_stability_corr) > abs(stops_stability_corr) * 1.5:
            primary_driver = "the **duration of stops**"
            primary_driver_is_duration = True
        else:
            primary_driver = "both the **frequency and duration of stops**"
        corr_insight = (f"This analysis suggests that <strong>{primary_driver}</strong> has the strongest impact on overall stability.")

    def format_period(period_value, level):
        if isinstance(period_value, (pd.Timestamp, pd.Period, pd.Timedelta)):
            return pd.to_datetime(period_value).strftime('%A, %b %d')
        if level == "Monthly": return f"Week {period_value}"
        if "Daily" in level: return f"{period_value}:00"
        return str(period_value)

    example_insight = ""
    if primary_driver_is_frequency:
        row = analysis_df_clean.loc[analysis_df_clean['stops'].idxmax()]
        example_insight = f"For example, the period with the most interruptions was <strong>{format_period(row['period'], analysis_level)}</strong>, which recorded <strong>{int(row['stops'])} stops</strong>."
    elif primary_driver_is_duration:
        row = analysis_df_clean.loc[analysis_df_clean['mttr'].idxmax()]
        example_insight = f"The period with the longest downtimes was <strong>{format_period(row['period'], analysis_level)}</strong>, where the average repair time was <strong>{row['mttr']:.1f} minutes</strong>."
    else:
        if not analysis_df_clean['mttr'].empty:
            row = analysis_df_clean.loc[analysis_df_clean['mttr'].idxmax()]
            example_insight = f"As an example, <strong>{format_period(row['period'], analysis_level)}</strong> experienced prolonged downtimes with an average repair time of <strong>{row['mttr']:.1f} minutes</strong>."

    return f"<div style='line-height: 1.6;'><p>{corr_insight}</p><p>{example_insight}</p></div>"


# ==============================================================================
# --- 5. EXCEL EXPORT MODULE ---
# ==============================================================================

def prepare_and_generate_run_based_excel(df_for_export, tolerance, downtime_gap_tolerance,
                                          run_interval_hours, tool_id_selection, startup_shots_count):
    try:
        base_calc = RunRateCalculator(df_for_export, tolerance, downtime_gap_tolerance,
                                      'aggregate', run_interval_hours, startup_shots_count)
        df_processed = base_calc.results.get("processed_df", pd.DataFrame())

        if df_processed.empty or 'run_id' not in df_processed.columns:
            st.error("Processing failed for Excel export.")
            return BytesIO().getvalue()

        mask_first = df_processed['tool_id'] != df_processed['tool_id'].shift(1)
        is_new_run = df_processed['time_diff_sec'] > (run_interval_hours * 3600)
        df_processed.loc[mask_first | is_new_run, 'stop_flag'] = 0
        df_processed['prev_stop_flag'] = (
            df_processed.groupby('tool_id')['stop_flag'].shift(1, fill_value=0)
        )
        df_processed['stop_event'] = (
            (df_processed['stop_flag'] == 1) & (df_processed['prev_stop_flag'] == 0)
        )
        df_processed['run_group'] = df_processed['stop_event'].cumsum()

        all_runs_data = {}
        desired_columns_base = [
            'SUPPLIER_NAME', 'tool_id', 'SESSION ID', 'shot_time',
            'APPROVED_CT', 'approved_ct', 'ACTUAL CT',
            'time_diff_sec', 'startup_flag', 'startup_event', 'stop_flag', 'stop_event', 'run_group'
        ]
        formula_columns = ['CUMULATIVE COUNT', 'RUN DURATION', 'TIME BUCKET']

        for run_id_val, df_run_raw in df_processed.groupby('run_id'):
            try:
                df_run_raw = df_run_raw.reset_index(drop=True)
                if df_run_raw.empty:
                    continue

                m = _run_metrics_from_processed(df_run_raw)
                run_results = {
                    'equipment_code': (df_run_raw['tool_id'].iloc[0] if 'tool_id' in df_run_raw.columns else tool_id_selection),
                    'start_time': m['start'],
                    'end_time': m['end'],
                    'mode_ct': m['mode_ct'],
                    'lower_limit': (df_run_raw['lower_limit'].iloc[0] if 'lower_limit' in df_run_raw.columns else 0),
                    'upper_limit': (df_run_raw['upper_limit'].iloc[0] if 'upper_limit' in df_run_raw.columns else np.inf),
                    'production_run_sec': m['duration'],
                    'total_runtime_sec': m['duration'],
                    'production_time_sec': m['prod_sec'],
                    'tot_down_time_sec': m['down_sec'],
                    'downtime_sec': m['down_sec'],
                    'mttr_min': ((m['down_sec'] / 60 / m['tot_stops']) if m['tot_stops'] > 0 else 0),
                    'mtbf_min': ((m['prod_sec'] / 60 / m['tot_stops']) if m['tot_stops'] > 0 else (m['prod_sec'] / 60)),
                    'avg_cycle_time_sec': (m['prod_sec'] / m['non_stop_shots'] if m['non_stop_shots'] > 0 else 0),
                }

                first_stop_idx = df_run_raw[df_run_raw['stop_event'] == True].index.min()
                if pd.isna(first_stop_idx):
                    time_to_first_dt = m['prod_sec']
                elif first_stop_idx == 0:
                    time_to_first_dt = 0
                else:
                    time_to_first_dt = df_run_raw.loc[:first_stop_idx - 1, 'adj_ct_sec'].sum()
                run_results['time_to_first_dt_min'] = time_to_first_dt / 60

                export_df = df_run_raw.copy()
                export_df['first_shot_time_diff'] = (export_df['time_diff_sec'].iloc[0] if not export_df.empty else 0)
                run_results['first_shot_time_diff'] = (export_df['first_shot_time_diff'].iloc[0] if not export_df.empty else 0)

                export_df['Shot Sequence'] = range(1, len(export_df) + 1)
                for col in formula_columns:
                    if col not in export_df.columns:
                        export_df[col] = ''

                cols_to_keep = [col for col in desired_columns_base if col in export_df.columns]
                cols_to_keep_final = (cols_to_keep + [col for col in formula_columns if col in export_df.columns])
                if 'Shot Sequence' in export_df.columns:
                    cols_to_keep_final.append('Shot Sequence')

                final_export_df = export_df[list(dict.fromkeys(cols_to_keep_final))].rename(
                    columns={
                        'tool_id': 'EQUIPMENT_CODE',
                        'shot_time': 'LOCAL_SHOT_TIME',
                        'time_diff_sec': 'TIME DIFF SEC',
                        'startup_flag': 'START-UP FLAG',
                        'startup_event': 'START-UP EVENT',
                        'stop_flag': 'STOP',
                        'stop_event': 'STOP EVENT'
                    }
                )

                final_desired_renamed = [
                    'SUPPLIER_NAME', 'EQUIPMENT_CODE', 'SESSION ID', 'Shot Sequence',
                    'LOCAL_SHOT_TIME', 'APPROVED_CT', 'approved_ct', 'ACTUAL CT',
                    'TIME DIFF SEC', 'START-UP FLAG', 'START-UP EVENT', 'STOP', 'STOP EVENT', 'run_group',
                    'CUMULATIVE COUNT', 'RUN DURATION', 'TIME BUCKET'
                ]
                for col in final_desired_renamed:
                    if col not in final_export_df.columns:
                        final_export_df[col] = ''
                final_export_df = final_export_df[
                    [col for col in final_desired_renamed if col in final_export_df.columns]
                ]

                run_results['processed_df'] = final_export_df
                all_runs_data[run_id_val] = run_results

            except Exception as e:
                st.warning(f"Could not process Run ID {run_id_val} for Excel: {e}")
                continue

        if not all_runs_data:
            st.error("No valid runs were processed for the Excel export.")
            return BytesIO().getvalue()

        return generate_excel_report(all_runs_data, tolerance)

    except Exception as e:
        st.error(f"Error preparing data for run-based Excel export: {e}")
        return BytesIO().getvalue()


def generate_excel_report(all_runs_data, tolerance):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        workbook = writer.book

        header_format = workbook.add_format({'bold': True, 'bg_color': '#002060', 'font_color': 'white', 'align': 'center', 'valign': 'vcenter', 'border': 1})
        sub_header_format = workbook.add_format({'bold': True, 'bg_color': '#C5D9F1', 'border': 1})
        label_format = workbook.add_format({'bold': True, 'align': 'left'})
        percent_format = workbook.add_format({'num_format': '0.0%', 'border': 1})
        time_format = workbook.add_format({'num_format': '[h]:mm:ss', 'border': 1})
        mins_format = workbook.add_format({'num_format': '0.00 "min"', 'border': 1})
        secs_format = workbook.add_format({'num_format': '0.00 "sec"', 'border': 1})
        data_format = workbook.add_format({'border': 1})
        datetime_format = workbook.add_format({'num_format': 'yyyy-mm-dd hh:mm:ss', 'border': 1})
        error_format = workbook.add_format({'bold': True, 'font_color': 'red'})

        for run_id, data in all_runs_data.items():
            ws = workbook.add_worksheet(f"Run_{run_id:03d}")
            df_run = data['processed_df'].copy()
            start_row = 19

            col_map = {name: xlsxwriter.utility.xl_col_to_name(i) for i, name in enumerate(df_run.columns)}
            shot_time_col_dyn = col_map.get('LOCAL_SHOT_TIME')
            stop_col = col_map.get('STOP')
            stop_event_col = col_map.get('STOP EVENT')
            time_bucket_col = col_map.get('TIME BUCKET')
            cum_count_col_dyn = col_map.get('CUMULATIVE COUNT')
            run_dur_col_dyn = col_map.get('RUN DURATION')
            bucket_col_dyn = col_map.get('TIME BUCKET')
            time_diff_col_dyn = col_map.get('TIME DIFF SEC')
            first_col_for_count = shot_time_col_dyn if shot_time_col_dyn else 'A'

            data_cols_count = len(df_run.columns)
            helper_col_letter = xlsxwriter.utility.xl_col_to_name(data_cols_count)
            ws.set_column(f'{helper_col_letter}:{helper_col_letter}', None, None, {'hidden': True})

            analysis_start_col_idx = data_cols_count + 2
            analysis_col_1 = xlsxwriter.utility.xl_col_to_name(analysis_start_col_idx)
            analysis_col_2 = xlsxwriter.utility.xl_col_to_name(analysis_start_col_idx + 1)
            analysis_col_3 = xlsxwriter.utility.xl_col_to_name(analysis_start_col_idx + 2)

            essential_cols = {
                'STOP': stop_col, 'STOP EVENT': stop_event_col,
                'TIME DIFF SEC': time_diff_col_dyn, 'CUMULATIVE COUNT': cum_count_col_dyn,
                'RUN DURATION': run_dur_col_dyn, 'TIME BUCKET': bucket_col_dyn,
                'LOCAL_SHOT_TIME': shot_time_col_dyn
            }
            missing_cols = [name for name, letter in essential_cols.items() if not letter]
            if missing_cols:
                ws.write('A5', f"Error: Missing columns: {', '.join(missing_cols)}", error_format)
            table_formulas_ok = not missing_cols

            ws.merge_range('A1:B1', data['equipment_code'], header_format)
            ws.write('A2', 'Date', label_format)
            ws.write('B2', f"{data['start_time']:%Y-%m-%d} to {data['end_time']:%Y-%m-%d}")
            ws.write('A3', 'Method', label_format)
            ws.write('B3', 'Every Shot')

            ws.write('E1', 'Mode CT', sub_header_format)
            mode_ct_val = data.get('mode_ct', 0)
            ws.write('E2', mode_ct_val if isinstance(mode_ct_val, (int, float)) else 0, secs_format)

            ws.write('F1', 'Outside L1', sub_header_format)
            ws.write('G1', 'Outside L2', sub_header_format)
            ws.write('H1', 'IDLE', sub_header_format)
            ws.write('F2', 'Lower Limit', label_format)
            ws.write('G2', 'Upper Limit', label_format)
            ws.write('H2', 'Stops', label_format)

            lower_limit_val = data.get('lower_limit')
            upper_limit_val = data.get('upper_limit')
            ws.write('F3', lower_limit_val if lower_limit_val is not None else 'N/A', secs_format)
            ws.write('G3', upper_limit_val if upper_limit_val is not None else 'N/A', secs_format)

            if stop_col:
                ws.write_formula('H3', f"=SUM({stop_col}{start_row}:{stop_col}{start_row + len(df_run) - 1})", sub_header_format)
            else:
                ws.write('H3', 'N/A', sub_header_format)

            ws.write('K1', 'Total Shot Count', label_format)
            ws.write('L1', 'Normal Shot Count', label_format)
            ws.write_formula('K2', f"=COUNTA({first_col_for_count}{start_row}:{first_col_for_count}{start_row + len(df_run) - 1})", sub_header_format)
            ws.write_formula('L2', "=K2-H3", sub_header_format)

            ws.write('K4', 'Efficiency', label_format)
            ws.write('L4', 'Stop Events', label_format)
            ws.write_formula('K5', "=L2/K2", percent_format)
            if stop_event_col:
                ws.write_formula('L5', f"=SUM({stop_event_col}{start_row}:{stop_event_col}{start_row + len(df_run) - 1})", sub_header_format)
            else:
                ws.write('L5', 'N/A', sub_header_format)

            ws.write('F5', 'Tot Run Time (Calc)', label_format)
            ws.write('G5', 'Tot Down Time', label_format)
            ws.write('H5', 'Tot Prod Time', label_format)

            downtime_to_write = data.get('tot_down_time_sec', 0)
            if not isinstance(downtime_to_write, (int, float)): downtime_to_write = 0

            ws.write('F6', data.get('production_run_sec', 0) / 86400, time_format)
            ws.write('G6', downtime_to_write / 86400, time_format)
            ws.write('H6', data.get('production_time_sec', 0) / 86400, time_format)

            ws.write('F4', '', label_format)
            ws.write('G4', 'Down %', label_format)
            ws.write('H4', 'Prod %', label_format)
            ws.write('F7', '', data_format)
            ws.write_formula('G7', "=IFERROR(G6/F6, 0)", percent_format)
            ws.write_formula('H7', "=IFERROR(H6/F6, 0)", percent_format)

            ws.merge_range('K8:L8', 'Reliability Metrics', header_format)
            ws.write('K9', 'MTTR (Avg)', label_format)
            ws.write('L9', data.get('mttr_min', 0), mins_format)
            ws.write('K10', 'MTBF (Avg)', label_format)
            ws.write('L10', data.get('mtbf_min', 0), mins_format)
            ws.write('K11', 'Time to First DT', label_format)
            ws.write('L11', data.get('time_to_first_dt_min', 0), mins_format)
            ws.write('K12', 'Avg Cycle Time', label_format)
            ws.write('L12', data.get('avg_cycle_time_sec', 0), secs_format)

            ws.merge_range(f'{analysis_col_1}14:{analysis_col_3}14', 'Time Bucket Analysis', header_format)
            ws.write(f'{analysis_col_1}15', 'Bucket', sub_header_format)
            ws.write(f'{analysis_col_2}15', 'Duration Range', sub_header_format)
            ws.write(f'{analysis_col_3}15', 'Events Count', sub_header_format)
            max_bucket = 20
            for i in range(1, max_bucket + 1):
                ws.write(f'{analysis_col_1}{15+i}', i, sub_header_format)
                ws.write(f'{analysis_col_2}{15+i}', f"{(i-1)*20} - {i*20} min", sub_header_format)
                if time_bucket_col:
                    ws.write_formula(f'{analysis_col_3}{15+i}', f'=COUNTIF({bucket_col_dyn}{start_row}:{bucket_col_dyn}{start_row + len(df_run) - 1},{i})', sub_header_format)
                else:
                    ws.write(f'{analysis_col_3}{15+i}', 'N/A', sub_header_format)
            ws.write(f'{analysis_col_2}{16+max_bucket}', 'Grand Total', sub_header_format)
            ws.write_formula(f'{analysis_col_3}{16+max_bucket}', f"=SUM({analysis_col_3}16:{analysis_col_3}{15+max_bucket})", sub_header_format)

            ws.write_row('A18', df_run.columns, header_format)

            df_run_nan_filled = df_run.fillna(np.nan)
            for i, row_values in enumerate(df_run_nan_filled.itertuples(index=False)):
                current_row_excel_idx = start_row + i - 1
                for c_idx, value in enumerate(row_values):
                    col_name = df_run.columns[c_idx]
                    if col_name in ['CUMULATIVE COUNT', 'RUN DURATION', 'TIME BUCKET', 'TIME DIFF SEC']:
                        continue

                    cell_format = data_format
                    if col_name in ['STOP', 'START-UP FLAG']:
                        ws.write_number(current_row_excel_idx, c_idx, int(value) if pd.notna(value) else 0, cell_format)
                    elif col_name in ['STOP EVENT', 'START-UP EVENT']:
                        ws.write_number(current_row_excel_idx, c_idx, 1 if value is True else 0, cell_format)
                    elif isinstance(value, pd.Timestamp):
                        if pd.notna(value):
                            value_no_tz = (value.tz_localize(None) if value.tzinfo is not None else value)
                            ws.write_datetime(current_row_excel_idx, c_idx, value_no_tz, datetime_format)
                        else:
                            ws.write_blank(current_row_excel_idx, c_idx, None, cell_format)
                    elif isinstance(value, (int, float, np.number)):
                        if col_name in ['ACTUAL CT', 'adj_ct_sec']: cell_format = secs_format
                        if pd.notna(value) and np.isfinite(value):
                            ws.write_number(current_row_excel_idx, c_idx, value, cell_format)
                        else:
                            ws.write_blank(current_row_excel_idx, c_idx, None, cell_format)
                    elif pd.isna(value):
                        ws.write_blank(current_row_excel_idx, c_idx, None, cell_format)
                    else:
                        ws.write_string(current_row_excel_idx, c_idx, str(value), cell_format)

            if table_formulas_ok:
                time_diff_col_idx = df_run.columns.get_loc('TIME DIFF SEC')
                cum_count_col_idx = df_run.columns.get_loc('CUMULATIVE COUNT')
                run_dur_col_idx = df_run.columns.get_loc('RUN DURATION')
                bucket_col_idx = df_run.columns.get_loc('TIME BUCKET')

                for i in range(len(df_run)):
                    row_num = start_row + i
                    prev_row = row_num - 1
                    current_row_zero_idx = start_row + i - 1

                    if i == 0:
                        first_diff_val = data.get('first_shot_time_diff', 0)
                        ws.write_number(current_row_zero_idx, time_diff_col_idx, first_diff_val, secs_format)
                    else:
                        formula = f'=IFERROR(({shot_time_col_dyn}{row_num}-{shot_time_col_dyn}{prev_row})*86400, 0)'
                        ws.write_formula(current_row_zero_idx, time_diff_col_idx, formula, secs_format)

                    if i == 0:
                        helper_formula = f'=IF({stop_col}{row_num}=0, {time_diff_col_dyn}{row_num}, 0)'
                    else:
                        helper_formula = f'=IF({stop_event_col}{row_num}=1, 0, IF({stop_col}{row_num}=0, {helper_col_letter}{prev_row}+{time_diff_col_dyn}{row_num}, {helper_col_letter}{prev_row}))'
                    ws.write_formula(current_row_zero_idx, data_cols_count, helper_formula)

                    cum_count_formula = f'=COUNTIF(${stop_event_col}${start_row}:${stop_event_col}{row_num},1)&"/"&IF({stop_event_col}{row_num}=1,"0 sec",TEXT({helper_col_letter}{row_num}/86400,"[h]:mm:ss"))'
                    ws.write_formula(current_row_zero_idx, cum_count_col_idx, cum_count_formula, data_format)

                    run_dur_formula = f'=IF({stop_event_col}{row_num}=1, IF({row_num}>{start_row}, {helper_col_letter}{prev_row}/86400, 0), "")'
                    ws.write_formula(current_row_zero_idx, run_dur_col_idx, run_dur_formula, time_format)

                    time_bucket_formula = f'=IF({stop_event_col}{row_num}=1, IF({row_num}>{start_row}, IFERROR(FLOOR({helper_col_letter}{prev_row}/60/20,1)+1, ""), ""), "")'
                    ws.write_formula(current_row_zero_idx, bucket_col_idx, time_bucket_formula, data_format)
            else:
                for col_dyn in [cum_count_col_dyn, time_diff_col_dyn, run_dur_col_dyn, bucket_col_dyn]:
                    if col_dyn: ws.write(f'{col_dyn}{start_row}', "Formula Error", error_format)

            for i, col_name in enumerate(df_run.columns):
                if col_name == "SESSION ID":
                    ws.set_column(i, i, None, None, {'hidden': True})
                    continue
                try:
                    max_len_data = df_run[col_name].astype(str).map(len).max()
                    max_len_data = 0 if pd.isna(max_len_data) else int(max_len_data)
                    width = max(len(str(col_name)), max_len_data)
                    ws.set_column(i, i, min(width + 2, 40))
                except Exception:
                    ws.set_column(i, i, len(str(col_name)) + 2)

    return output.getvalue()


# ==============================================================================
# --- 6. RISK ANALYSIS MODULE ---
# ==============================================================================

def calculate_risk_scores(df_all, run_interval_hours=8, min_shots_filter=1, tolerance=0.05, downtime_gap_tolerance=2.0, startup_shots_count=0):
    if df_all.empty or 'tool_id' not in df_all.columns:
        return pd.DataFrame()

    initial_metrics = []

    for tool_id, df_tool in df_all.groupby('tool_id'):
        df_tool = df_tool.sort_values(['shot_time', 'ACTUAL CT'])
        if df_tool.empty:
            continue

        max_date = df_tool['shot_time'].max()
        cutoff_date = max_date - timedelta(weeks=4)
        df_period = df_tool[df_tool['shot_time'] >= cutoff_date].copy()
        if df_period.empty:
            continue

        calc = RunRateCalculator(df_period, tolerance, downtime_gap_tolerance, 'aggregate', run_interval_hours, startup_shots_count)
        res = calc.results
        df_processed = res.get('processed_df')
        if df_processed is None or df_processed.empty:
            continue

        run_summary_df = calculate_run_summaries(df_processed, tolerance, downtime_gap_tolerance, run_interval_hours, pre_processed=True, startup_shots_count=startup_shots_count)
        if run_summary_df.empty:
            continue

        if 'total_shots' in run_summary_df.columns:
            run_summary_df = run_summary_df[run_summary_df['total_shots'] >= min_shots_filter]
        if run_summary_df.empty:
            continue

        total_runtime_sec = run_summary_df['total_runtime_sec'].sum()
        production_time_sec = run_summary_df['production_time_sec'].sum()
        downtime_sec = run_summary_df['downtime_sec'].sum()
        stop_events = run_summary_df['stops'].sum()

        res_stability = ((production_time_sec / total_runtime_sec * 100) if total_runtime_sec > 0 else 100.0)
        res_mttr = (downtime_sec / 60 / stop_events) if stop_events > 0 else 0
        res_mtbf = ((production_time_sec / 60 / stop_events) if stop_events > 0 else (production_time_sec / 60))

        weekly_stats = []
        df_processed = df_processed.copy()
        iso = df_processed['shot_time'].dt.isocalendar()
        df_processed['week'] = iso['week']
        df_processed['iso_year'] = iso['year']
        df_processed['week_key'] = df_processed['iso_year'].astype(str) + '-W' + df_processed['week'].astype(str).str.zfill(2)

        sorted_weeks = sorted(
            [(g_df['shot_time'].min(), w_key, g_df)
             for w_key, g_df in df_processed.groupby('week_key') if not g_df.empty],
            key=lambda x: x[0]
        )

        for _, week_key, df_week in sorted_weeks:
            weekly_run_summary = calculate_run_summaries(df_week.copy(), tolerance, downtime_gap_tolerance,
                                                          run_interval_hours, pre_processed=True, startup_shots_count=startup_shots_count)
            if not weekly_run_summary.empty:
                if 'total_shots' in weekly_run_summary.columns:
                    weekly_run_summary = weekly_run_summary[weekly_run_summary['total_shots'] >= min_shots_filter]
                if not weekly_run_summary.empty:
                    w_tot_runtime = weekly_run_summary['total_runtime_sec'].sum()
                    w_prod_time = weekly_run_summary['production_time_sec'].sum()
                    w_stability = ((w_prod_time / w_tot_runtime * 100) if w_tot_runtime > 0 else 100.0)
                    weekly_stats.append({'week': week_key, 'stability': w_stability})

        weekly_stabilities_df = pd.DataFrame(weekly_stats)
        weekly_stabilities = (weekly_stabilities_df['stability'].tolist() if not weekly_stabilities_df.empty else [])

        trend = "Stable"
        if (len(weekly_stabilities) > 1 and weekly_stabilities[-1] < weekly_stabilities[0] * 0.95):
            trend = "Declining"

        initial_metrics.append({
            'Tool ID': tool_id,
            'Stability': res_stability,
            'MTTR': res_mttr,
            'MTBF': res_mtbf,
            'Weekly Stability': ' → '.join([f'{s:.0f}%' for s in weekly_stabilities]),
            'Trend': trend,
            'Analysis Period': f"{cutoff_date.strftime('%Y-%m-%d')} to {max_date.strftime('%Y-%m-%d')}",
        })

    if not initial_metrics:
        return pd.DataFrame()

    metrics_df = pd.DataFrame(initial_metrics)
    overall_mttr_mean = metrics_df['MTTR'].mean()
    overall_mtbf_mean = metrics_df['MTBF'].mean()

    final_risk_data = []
    for _, row in metrics_df.iterrows():
        risk_score = row['Stability']
        if row['Trend'] == "Declining":
            risk_score -= 20

        primary_factor = "Stable"
        details = f"Overall stability is {row['Stability']:.1f}%."

        if row['Trend'] == "Declining":
            primary_factor = "Declining Trend"
            details = "Declining stability"
        elif (row['Stability'] < 70 and overall_mttr_mean > 0 and row['MTTR'] > (overall_mttr_mean * 1.2)):
            primary_factor = "High MTTR"
            details = f"Avg stop duration (MTTR) of {row['MTTR']:.1f} min is high."
        elif (row['Stability'] < 70 and overall_mtbf_mean > 0 and row['MTBF'] < (overall_mtbf_mean * 0.8)):
            primary_factor = "Frequent Stops"
            details = f"Frequent stops (MTBF of {row['MTBF']:.1f} min)."
        elif row['Stability'] <= 50:
            primary_factor = "Critical Stability"
            details = f"Overall stability is critical ({row['Stability']:.1f}%)."
        elif row['Stability'] <= 70:
            primary_factor = "Moderate Stability"
            details = f"Stability is below target ({row['Stability']:.1f}%)."

        final_risk_data.append({
            'Tool ID': row['Tool ID'],
            'Analysis Period': row['Analysis Period'],
            'Risk Score': max(0, risk_score),
            'Primary Risk Factor': primary_factor,
            'Weekly Stability': row['Weekly Stability'],
            'Details': details,
        })

    if not final_risk_data:
        return pd.DataFrame()

    return (pd.DataFrame(final_risk_data)
              .sort_values('Risk Score', ascending=True)
              .reset_index(drop=True))


# ==============================================================================
# --- 7. WEEKLY COMPARISON REPORT GENERATOR ---
# ==============================================================================

def generate_weekly_comparison_pptx(df_weekly: pd.DataFrame, tool_id: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(hex_str):
        h = hex_str.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _fmt(val, fmt):
        try:
            if pd.isna(val): return "—"
            return fmt.format(val)
        except Exception:
            return str(val)

    def _delta_str(curr, prev, higher_is_better=True):
        try:
            if pd.isna(prev) or prev == 0: return "", None
            pct = (curr - prev) / abs(prev) * 100
            sign = "+" if pct >= 0 else ""
            label = f"{{{sign}{pct:.1f}%}}"
            good  = pct >= 0 if higher_is_better else pct <= 0
            color = "2E7D32" if good else "C62828"
            return label, color
        except Exception:
            return "", None

    METRICS = [
        ("Stability Index",     "Stability Index (%)",  "{:.1f}%",  True),
        ("Run Rate Efficiency", "Efficiency (%)",        "{:.1f}%",  True),
        ("MTTR",                "MTTR (min)",            "{:.1f} min", False),
        ("MTBF",                "MTBF (min)",            "{:.1f} min", True),
        ("Production Time",     "Production Time (h)",  "{:.1f} h",  True),
        ("Total Shots",         "Total Shots",           "{:,.0f}",   True),
        ("Normal Shots",        "Normal Shots",          "{:,.0f}",   True),
        ("Stop Events",         "Stop Events",           "{:.0f}",    False),
    ]

    df = df_weekly.reset_index(drop=True)
    periods = df["Week"].tolist()
    n_weeks = len(periods)

    totals = {}
    for _, col, fmt, _ in METRICS:
        if col not in df.columns:
            totals[col] = None
            continue
        if col in ("Total Shots", "Normal Shots", "Stop Events"):
            totals[col] = df[col].sum()
        else:
            totals[col] = df[col].mean()

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("F4F6F9")

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Tooling Performance: {tool_id}  —  Weekly Comparison"
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = rgb("002060")
    run.font.name  = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(1, Inches(0.3), Inches(0.92), Inches(12.7), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = rgb("002060")
    line.line.fill.background()

    left   = Inches(0.3)
    top    = Inches(1.05)
    width  = Inches(12.7)
    height = Inches(5.7)

    n_cols   = 1 + n_weeks + 1
    n_rows   = 1 + len(METRICS)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    kpi_w    = Inches(1.9)
    data_w   = (Inches(12.7) - kpi_w) / (n_cols - 1)
    tbl.columns[0].width = kpi_w
    for c in range(1, n_cols):
        tbl.columns[c].width = int(data_w)

    hdr_h  = Inches(0.52)
    row_h  = (Inches(5.7) - hdr_h) / len(METRICS)
    tbl.rows[0].height = int(hdr_h)
    for r in range(1, n_rows):
        tbl.rows[r].height = int(row_h)

    HDR_BG   = rgb("002060")
    HDR_FG   = rgb("FFFFFF")
    ALT_BG   = rgb("EEF2F7")
    NORM_BG  = rgb("FFFFFF")
    KPI_FG   = rgb("1A237E")
    VAL_FG   = rgb("212121")
    TOT_BG   = rgb("D9E1F2")
    TOT_FG   = rgb("002060")

    def _cell_set(cell, text, fg, bg, bold=False, size=10, align=PP_ALIGN.CENTER, italic=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf2 = cell.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = align
        for run2 in p2.runs: run2.text = ""
        if p2.runs: r2 = p2.runs[0]
        else: r2 = p2.add_run()
        r2.text = text
        r2.font.size  = Pt(size)
        r2.font.bold  = bold
        r2.font.italic = italic
        r2.font.color.rgb = fg
        r2.font.name  = "Calibri"

    def _cell_rich(cell, main_text, delta_text, delta_color_hex, bg, main_fg, main_size=10, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf2 = cell.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        for _ in list(p2.runs): pass
        r_main = p2.add_run()
        r_main.text = main_text + "  "
        r_main.font.size  = Pt(main_size)
        r_main.font.bold  = bold
        r_main.font.color.rgb = main_fg
        r_main.font.name  = "Calibri"
        if delta_text:
            r_delta = p2.add_run()
            r_delta.text = delta_text
            r_delta.font.size  = Pt(8)
            r_delta.font.bold  = False
            r_delta.font.color.rgb = rgb(delta_color_hex)
            r_delta.font.name  = "Calibri"

    _cell_set(tbl.cell(0, 0), "KPI", HDR_FG, HDR_BG, bold=True, size=11, align=PP_ALIGN.LEFT)
    for wi, period in enumerate(periods):
        _cell_set(tbl.cell(0, wi + 1), str(period), HDR_FG, HDR_BG, bold=True, size=11)
    _cell_set(tbl.cell(0, n_cols - 1), "Total / Avg", HDR_FG, HDR_BG, bold=True, size=11)

    for ri, (label, col, fmt, hib) in enumerate(METRICS):
        row_idx = ri + 1
        row_bg  = ALT_BG if ri % 2 == 0 else NORM_BG

        _cell_set(tbl.cell(row_idx, 0), label, KPI_FG, row_bg, bold=True, size=10, align=PP_ALIGN.LEFT)

        for wi in range(n_weeks):
            curr_val = df.iloc[wi][col] if col in df.columns else None
            prev_val = df.iloc[wi - 1][col] if (wi > 0 and col in df.columns) else None

            val_str   = _fmt(curr_val, fmt)
            delta_str, delta_color = ("", None)
            if wi > 0 and curr_val is not None and prev_val is not None:
                delta_str, delta_color = _delta_str(curr_val, prev_val, hib)

            cell = tbl.cell(row_idx, wi + 1)
            if delta_str and delta_color:
                _cell_rich(cell, val_str, delta_str, delta_color, row_bg, VAL_FG, main_size=10, bold=(wi == 0))
            else:
                _cell_set(cell, val_str, VAL_FG, row_bg, bold=(wi == 0), size=10)

        tot_val  = totals.get(col)
        tot_str  = _fmt(tot_val, fmt) if tot_val is not None else "—"
        _cell_set(tbl.cell(row_idx, n_cols - 1), tot_str, TOT_FG, TOT_BG, bold=True, size=10)

    note_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.4))
    tf3 = note_box.text_frame
    p3  = tf3.paragraphs[0]
    r3  = p3.add_run()
    r3.text  = f"Generated by Run Rate Analysis v3.50  |  Tool: {tool_id}  |  {pd.Timestamp.now().strftime('%d %b %Y')}"
    r3.font.size   = Pt(8)
    r3.font.italic = True
    r3.font.color.rgb = rgb("9E9E9E")
    r3.font.name   = "Calibri"
    p3.alignment = PP_ALIGN.LEFT

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()